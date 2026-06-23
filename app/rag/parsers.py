"""
=============================================================================
 Enterprise Level RAG: Layer 1 — Universal Document Parser
=============================================================================
 Supports: PDF, DOCX, XLSX, PPTX, CSV, TXT, Images, Video (subtitles)
 All parsing is 100% offline — no API calls, no cloud services.
=============================================================================
"""

import csv
import io
import os
import re
import shutil
import subprocess
import email
from email import policy
import urllib.request
from dataclasses import dataclass, field
from typing import List, Optional

# ---------------------------------------------------------------------------
# Optional imports — each format gracefully degrades if library missing
# ---------------------------------------------------------------------------
try:
    from docling.document_converter import DocumentConverter
    DOCLING_AVAILABLE = True
except ImportError as e:
    import traceback
    print(f"[Parsers] Docling import failed: {e}")
    traceback.print_exc()
    DOCLING_AVAILABLE = False

try:
    from app.rag.table_engine import (
        markdown_to_rich_table,
        annotate_section_title,
        extract_tables_pdfplumber,
        stitch_continuation_tables,
    )
    TABLE_ENGINE_AVAILABLE = True
except ImportError:
    TABLE_ENGINE_AVAILABLE = False



try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    import pysubs2
    SUBS_AVAILABLE = True
except ImportError:
    SUBS_AVAILABLE = False

try:
    import chardet
    CHARDET_AVAILABLE = True
except ImportError:
    CHARDET_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ---------------------------------------------------------------------------
# Data Models
# ---------------------------------------------------------------------------
@dataclass
class PageContent:
    """Content extracted from a single page / sheet / slide."""
    page_num: int
    text: str = ""
    tables: List[str] = field(default_factory=list)
    rich_tables: List[object] = field(default_factory=list)   # RichTable objects (table_engine)
    image_texts: List[str] = field(default_factory=list)
    image_bytes: List[bytes] = field(default_factory=list)
    content_type: str = "text"  # text, table, image_ocr, subtitle
    section_title: str = ""    # Most recent heading above this content


@dataclass
class ParseResult:
    """Complete result from parsing a file."""
    pages: List[PageContent] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)
    success: bool = True
    error: Optional[str] = None


# ---------------------------------------------------------------------------
# Supported Extensions
# ---------------------------------------------------------------------------
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".gif", ".webp"}
VIDEO_EXTENSIONS = {".mp4", ".avi", ".mkv", ".mov", ".wmv", ".flv"}
SUBTITLE_EXTENSIONS = {".srt", ".ass", ".ssa", ".vtt"}
TEXT_EXTENSIONS = {".txt", ".text", ".md", ".log", ".json", ".xml"}
CODE_EXTENSIONS = {".py", ".js", ".java", ".cpp", ".c", ".h", ".cs", ".go", ".rs", ".rb", ".php", ".sql", ".html", ".css", ".sh"}
EMAIL_EXTENSIONS = {".eml", ".msg"}
URL_EXTENSIONS = {".url", ".webloc"}
ARCHIVE_EXTENSIONS = {".zip", ".tar", ".gz", ".tgz", ".rar"}
MAX_ARCHIVE_EXTRACT_SIZE = 1024 * 1024 * 1024  # 1GB limit for zip bomb protection

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".csv",
    ".pptx", ".ppt"
} | IMAGE_EXTENSIONS | VIDEO_EXTENSIONS | SUBTITLE_EXTENSIONS | TEXT_EXTENSIONS | CODE_EXTENSIONS | EMAIL_EXTENSIONS | URL_EXTENSIONS | ARCHIVE_EXTENSIONS


def is_supported_file(file_path: str) -> bool:
    """Check if a file extension is supported."""
    ext = os.path.splitext(file_path)[1].lower()
    return ext in SUPPORTED_EXTENSIONS


def get_file_type(file_path: str) -> str:
    """Return a human-readable file type category."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return "pdf"
    elif ext in {".docx", ".doc"}:
        return "docx"
    elif ext in {".xlsx", ".xls"}:
        return "xlsx"
    elif ext in {".pptx", ".ppt"}:
        return "pptx"
    elif ext == ".csv":
        return "csv"
    elif ext in TEXT_EXTENSIONS:
        return "text"
    elif ext in CODE_EXTENSIONS:
        return "code"
    elif ext in EMAIL_EXTENSIONS:
        return "email"
    elif ext in URL_EXTENSIONS:
        return "url"
    elif ext in ARCHIVE_EXTENSIONS:
        return "archive"
    elif ext in IMAGE_EXTENSIONS:
        return "image"
    elif ext in VIDEO_EXTENSIONS:
        return "video"
    elif ext in SUBTITLE_EXTENSIONS:
        return "subtitle"
    return "unknown"


# ---------------------------------------------------------------------------
# Utility Functions
# ---------------------------------------------------------------------------
def _auto_rotate_image(img):
    """Detect image orientation via OSD and auto-rotate upright (solves PDF flip)."""
    if not OCR_AVAILABLE:
        return img
    try:
        osd = pytesseract.image_to_osd(img)
        match = re.search(r'(?<=Rotate: )\d+', osd)
        if match:
            angle = int(match.group(0))
            if angle != 0:
                # pytesseract's Rotate value is counter-clockwise angle to correct orientation
                img = img.rotate(angle, expand=True)
    except Exception:
        pass  # OSD fails if too little text is detected
    return img


def _enrich_mcq_text(text: str) -> str:
    """Detect MCQ checkmarks and explicitly mark the correct answer."""
    if not text:
        return text
    
    # Common symbols used for marking the correct MCQ option
    mcq_symbols = [r"✓", r"✔", r"☑", r"☒", r"◉", r"\[x\]", r"\[X\]", r"\(x\)", r"\(X\)"]
    pattern = re.compile(f"({'|'.join(mcq_symbols)})")
    
    enriched_lines = []
    for line in text.split('\n'):
        if pattern.search(line):
            # If the line has an MCQ tick, explicitly label it so the LLM understands it's the correct answer
            line = f"[CORRECT ANSWER] {line}"
        enriched_lines.append(line)
    return "\n".join(enriched_lines)





# ---------------------------------------------------------------------------
# Superscript/Subscript Normalization
# ---------------------------------------------------------------------------
# Unicode superscript/subscript characters that PDFs commonly emit in catalogue
# model numbers (e.g., EQL40200D³, ½, ¼). These MUST be normalized to plain
# ASCII so that keyword search for "EQL40200D" matches "EQL40200D3".
_SUPERSCRIPT_MAP = str.maketrans({
    '⁰': '0', '¹': '1', '²': '2', '³': '3', '⁴': '4',
    '⁵': '5', '⁶': '6', '⁷': '7', '⁸': '8', '⁹': '9',
    '₀': '0', '₁': '1', '₂': '2', '₃': '3', '₄': '4',
    '₅': '5', '₆': '6', '₇': '7', '₈': '8', '₉': '9',
    '½': '1/2', '¼': '1/4', '¾': '3/4',
    '⅛': '1/8', '⅜': '3/8', '⅝': '5/8', '⅞': '7/8',
    '⅓': '1/3', '⅔': '2/3',
    '®': '(R)', '™': '(TM)', '©': '(C)',
    '①': ' [Note 1]', '②': ' [Note 2]', '③': ' [Note 3]', '④': ' [Note 4]', '⑤': ' [Note 5]',
    '⑥': ' [Note 6]', '⑦': ' [Note 7]', '⑧': ' [Note 8]', '⑨': ' [Note 9]', '⑩': ' [Note 10]',
})


def _normalize_superscripts(text: str) -> str:
    """Normalize Unicode superscripts, subscripts, fractions, and circled numbers
    to plain ASCII equivalents for accurate keyword matching."""
    if not text:
        return text
    # Handle multi-character replacements first
    text = text.replace('¹⁄', '1/').replace('³⁄', '3/').replace('⁵⁄', '5/').replace('⁷⁄', '7/')
    text = text.translate(_SUPERSCRIPT_MAP)
    
    # --- RegEx Enhancement: Splitting Merged Columns ---
    # Detects when a Catalogue/Model Number is jammed against a Measurement (Amperage/Voltage)
    # e.g., "EQL40200D100A" -> "EQL40200D | 100A"
    # Matches: [Model ID 4-10 chars] followed optionally by space, then [Number][A|V|W]
    text = re.sub(r'\b([A-Z0-9]{4,12})\s*(\d{1,4}[AVW])\b', r'\1 | \2', text)
    
    # Detects Model Number jammed against a dimensional string or phase e.g. "SNC2448L1125 1 Phase"
    text = re.sub(r'\b([A-Z0-9]{4,12})\s*(\d{1}\s*Phase)\b', r'\1 | \2', text)
    
    return text


# ---------------------------------------------------------------------------
# Topic/Heading Tracker Heuristics
# ---------------------------------------------------------------------------





# ---------------------------------------------------------------------------
# Office/HTML Parser (Docling)
# ---------------------------------------------------------------------------
def _parse_docling(file_path: str) -> ParseResult:
    """Extract all content from PDF using IBM Docling.
    Captures: text (layout-aware), tables (perfect structure), images.
    Replaces older PyMuPDF + pdfplumber + vLLM pipeline.
    """
    if not DOCLING_AVAILABLE:
        return ParseResult(success=False, error="docling not installed")

    try:
        from docling.document_converter import DocumentConverter
        from app.rag.model_loader import get_optimal_device

        # Build converter with optimal options (v1 API) or without (v2 auto-detect)
        try:
            from docling.document_converter import PdfFormatOption
            from docling.datamodel.pipeline_options import PdfPipelineOptions, TableStructureOptions, OcrOptions
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.accelerator_options import AcceleratorOptions, AcceleratorDevice

            optimal = get_optimal_device()
            device_mapping = {
                "cuda": AcceleratorDevice.CUDA,
                "mps": AcceleratorDevice.MPS,
                "cpu": AcceleratorDevice.CPU
            }
            accel_device = device_mapping.get(optimal, AcceleratorDevice.AUTO)

            pipeline_options = PdfPipelineOptions()
            pipeline_options.accelerator_options = AcceleratorOptions(
                num_threads=8,
                device=accel_device
            )
            
            # --- MOST POWERFUL VERSION CONFIGURATION ---
            pipeline_options.do_ocr = True
            pipeline_options.generate_page_images = True
            pipeline_options.generate_picture_images = True
            pipeline_options.do_table_structure = True
            pipeline_options.table_structure_options = TableStructureOptions(mode="accurate")

            converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
                }
            )
        except (ImportError, AttributeError):
            # Docling v2+ API — auto-detect format, no explicit options needed
            converter = DocumentConverter()

        result = converter.convert(file_path)
        doc = result.document

        pages_dict = {}
        # Track last seen section heading so each table gets its owning title
        last_section_header: str = ""
        # Accumulate tables per page for stitching
        tables_by_page: dict = {}
        table_counter = 0

        # Iterate through all layout items detected by Docling
        for item, level in doc.iterate_items():
            page_no = 1
            if hasattr(item, "prov") and item.prov:
                page_no = item.prov[0].page_no

            if page_no not in pages_dict:
                pages_dict[page_no] = PageContent(
                    page_num=page_no, text="", tables=[],
                    rich_tables=[], image_texts=[], image_bytes=[]
                )

            page_content = pages_dict[page_no]
            item_type = type(item).__name__

            if item_type == "SectionHeaderItem":
                # Update the running section header tracker
                header_text = _normalize_superscripts(item.text)
                last_section_header = header_text
                page_content.text += f"\n## {header_text}\n\n"
                page_content.section_title = header_text

            elif item_type == "TableItem":
                table_md = item.export_to_markdown(doc)
                table_md = _normalize_superscripts(table_md)
                # Prepend the owning section title to the table label
                section_label = f" (section: {last_section_header})" if last_section_header else ""
                table_text = f"[TABLE - Page {page_no}{section_label}]\n{table_md}"
                page_content.tables.append(table_text)

                # Build a RichTable from the Docling markdown
                if TABLE_ENGINE_AVAILABLE:
                    table_id = f"p{page_no}_t{table_counter}"
                    rich = markdown_to_rich_table(
                        table_md,
                        table_id=table_id,
                        section_title=last_section_header,
                        page_no=page_no,
                    )
                    if rich:
                        page_content.rich_tables.append(rich)
                        if page_no not in tables_by_page:
                            tables_by_page[page_no] = []
                        tables_by_page[page_no].append(rich)
                    table_counter += 1

            elif item_type == "PictureItem":
                if hasattr(item, "caption") and item.caption:
                    caption_text = _normalize_superscripts(item.caption.text)
                    page_content.image_texts.append(f"[IMAGE CAPTION - Page {page_no}]\n{caption_text}")

                try:
                    img = item.get_image(doc)
                    if img:
                        img_byte_arr = io.BytesIO()
                        img.save(img_byte_arr, format='PNG')
                        page_content.image_bytes.append(img_byte_arr.getvalue())
                except Exception as img_err:
                    print(f"[Docling] Failed to extract image on page {page_no}: {img_err}")

            elif item_type == "TextItem":
                clean_text = _enrich_mcq_text(item.text)
                clean_text = _normalize_superscripts(clean_text)
                page_content.text += clean_text + "\n\n"

        # Multi-page table stitching
        if TABLE_ENGINE_AVAILABLE and tables_by_page:
            try:
                stitched = stitch_continuation_tables(tables_by_page)
                # Sync stitched rich_tables back into pages_dict
                for pg, rich_list in stitched.items():
                    if pg in pages_dict:
                        pages_dict[pg].rich_tables = rich_list
            except Exception as stitch_err:
                print(f"[Docling] Table stitching failed: {stitch_err}")

        pages = [pages_dict[k] for k in sorted(pages_dict.keys())]

        doc_metadata = {
            "format": "pdf",
            "page_count": len(pages),
            "source": file_path,
        }

        return ParseResult(
            pages=pages,
            metadata=doc_metadata,
        )
    except Exception as e:
        return ParseResult(success=False, error=f"Docling PDF parse error: {e}")





# ---------------------------------------------------------------------------
# CSV Parser
# ---------------------------------------------------------------------------
def _parse_csv(file_path: str) -> ParseResult:
    """Parse CSV files into markdown table format."""
    try:
        # Detect encoding
        encoding = "utf-8"
        if CHARDET_AVAILABLE:
            with open(file_path, "rb") as f:
                raw = f.read(10000)
                detected = chardet.detect(raw)
                encoding = detected.get("encoding", "utf-8") or "utf-8"

        with open(file_path, "r", encoding=encoding, errors="replace") as f:
            reader = csv.reader(f)
            rows = []
            previous_row = None
            
            for row in reader:
                # Clean and normalize superscripts
                cleaned_row = [_normalize_superscripts(str(cell).strip()) if cell else "" for cell in row]
                
                # Conservative empty cell inheritance for CSVs
                if previous_row and len(cleaned_row) == len(previous_row):
                    if any(cleaned_row):
                        for i in range(len(cleaned_row)):
                            if not cleaned_row[i] and previous_row[i]:
                                cleaned_row[i] = previous_row[i]
                
                if any(cleaned_row):
                    rows.append(" | ".join(cleaned_row))
                    previous_row = cleaned_row

        if not rows:
            return ParseResult(
                pages=[PageContent(page_num=1, text="(Empty CSV file)")],
                metadata={"format": "csv", "page_count": 1, "source": file_path},
            )

        header = rows[0]
        col_count = len(rows[0].split(" | "))
        separator = " | ".join(["---"] * col_count)
        table_md = f"[CSV DATA]\n{header}\n{separator}\n" + "\n".join(rows[1:])

        # Split into pages of 100 rows for very large CSVs
        chunk_size = 100
        pages = []
        data_rows = rows[1:]
        if not data_rows:
            pages.append(PageContent(
                page_num=1,
                text=f"[CSV DATA]\n{header}\n{separator}\n(No data rows)",
            ))
        for i in range(0, len(data_rows), chunk_size):
            chunk = data_rows[i:i + chunk_size]
            chunk_table = f"[CSV DATA - Rows {i + 1} to {i + len(chunk)}]\n{header}\n{separator}\n" + "\n".join(chunk)
            pages.append(PageContent(
                page_num=(i // chunk_size) + 1,
                text=chunk_table,
                tables=[chunk_table],
                content_type="table",
            ))

        if not pages:
            pages.append(PageContent(page_num=1, text=table_md, tables=[table_md], content_type="table"))

        return ParseResult(
            pages=pages,
            metadata={
                "format": "csv",
                "page_count": len(pages),
                "total_rows": len(rows),
                "source": file_path,
            },
        )
    except Exception as e:
        return ParseResult(success=False, error=f"CSV parse error: {e}")


# ---------------------------------------------------------------------------
# Plain Text Parser
# ---------------------------------------------------------------------------
def _parse_text(file_path: str) -> ParseResult:
    """Parse plain text, markdown, log, JSON, XML files."""
    try:
        encoding = "utf-8"
        if CHARDET_AVAILABLE:
            with open(file_path, "rb") as f:
                raw = f.read(10000)
                detected = chardet.detect(raw)
                encoding = detected.get("encoding", "utf-8") or "utf-8"

        with open(file_path, "r", encoding=encoding, errors="replace") as f:
            content = f.read()
            
        content = _enrich_mcq_text(content)

        if not content.strip():
            return ParseResult(
                pages=[PageContent(page_num=1, text="(Empty file)")],
                metadata={"format": "text", "page_count": 1, "source": file_path},
            )

        # Split large text files into ~4000 char pages
        page_size = 4000
        pages = []
        lines = content.split("\n")
        current_page = []
        current_len = 0
        page_num = 1

        for line in lines:
            if current_len + len(line) > page_size and current_page:
                pages.append(PageContent(
                    page_num=page_num,
                    text="\n".join(current_page),
                ))
                page_num += 1
                current_page = []
                current_len = 0
            current_page.append(line)
            current_len += len(line) + 1

        if current_page:
            pages.append(PageContent(page_num=page_num, text="\n".join(current_page)))

        return ParseResult(
            pages=pages,
            metadata={
                "format": "text",
                "page_count": len(pages),
                "total_chars": len(content),
                "source": file_path,
            },
        )
    except Exception as e:
        return ParseResult(success=False, error=f"Text parse error: {e}")


# ---------------------------------------------------------------------------
# Image Parser (OCR)
# ---------------------------------------------------------------------------
def _parse_image(file_path: str) -> ParseResult:
    """Extract text from images using OCR. Also passes raw bytes for CLIP vision."""
    if not OCR_AVAILABLE:
        return ParseResult(success=False, error="pytesseract/Pillow not installed for image OCR")

    try:
        # Always read raw bytes first for CLIP vision pipeline
        with open(file_path, "rb") as fh:
            raw_bytes = fh.read()

        img = Image.open(io.BytesIO(raw_bytes))

        # Apply EXIF orientation before processing
        try:
            exif = img.getexif()
            if exif:
                orientation = exif.get(274)
                if orientation == 3:
                    img = img.rotate(180, expand=True)
                elif orientation == 6:
                    img = img.rotate(270, expand=True)
                elif orientation == 8:
                    img = img.rotate(90, expand=True)
        except Exception:
            pass

        # Pre-process for optimal OCR
        if img.mode != "L":
            img = img.convert("L")

        img = _auto_rotate_image(img)
        ocr_text = pytesseract.image_to_string(img).strip()
        ocr_text = _enrich_mcq_text(ocr_text)
        filename = os.path.basename(file_path)

        base = PageContent(
            page_num=1,
            image_texts=[ocr_text] if ocr_text else [],
            content_type="image_ocr",
        )

        if not ocr_text:
            base.text = f"[IMAGE: {filename}]\n(No text detected in image)"
            base.image_bytes = [raw_bytes]
            return ParseResult(
                pages=[base],
                metadata={"format": "image", "page_count": 1, "source": file_path},
            )

        base.text = f"[IMAGE OCR: {filename}]\n{ocr_text}"
        base.image_bytes = [raw_bytes]
        return ParseResult(
            pages=[base],
            metadata={
                "format": "image",
                "page_count": 1,
                "image_size": f"{img.width}x{img.height}",
                "source": file_path,
            },
        )
    except Exception as e:
        return ParseResult(success=False, error=f"Image parse error: {e}")


# ---------------------------------------------------------------------------
# Video Subtitle Parser
# ---------------------------------------------------------------------------
def _parse_video(file_path: str) -> ParseResult:
    """Extract embedded subtitles from video files using ffmpeg + pysubs2."""
    try:
        # First try to extract embedded subtitles using ffmpeg
        subtitle_text = _extract_subtitles_ffmpeg(file_path)

        if not subtitle_text:
            return ParseResult(
                pages=[PageContent(
                    page_num=1,
                    text=f"[VIDEO: {os.path.basename(file_path)}]\n(No embedded subtitles found)",
                    content_type="subtitle",
                )],
                metadata={"format": "video", "page_count": 1, "source": file_path},
            )

        # Split subtitles into manageable pages
        lines = subtitle_text.split("\n")
        page_size = 50  # lines per page
        pages = []
        for i in range(0, len(lines), page_size):
            chunk = lines[i:i + page_size]
            pages.append(PageContent(
                page_num=(i // page_size) + 1,
                text=f"[VIDEO SUBTITLES - Part {(i // page_size) + 1}]\n" + "\n".join(chunk),
                content_type="subtitle",
            ))

        return ParseResult(
            pages=pages,
            metadata={
                "format": "video",
                "page_count": len(pages),
                "subtitle_lines": len(lines),
                "source": file_path,
            },
        )
    except Exception as e:
        return ParseResult(success=False, error=f"Video parse error: {e}")


def _extract_subtitles_ffmpeg(video_path: str) -> str:
    """Extract embedded subtitle tracks from video using ffmpeg."""
    try:
        # Extract subtitle to SRT format using ffmpeg
        result = subprocess.run(
            [
                "ffmpeg", "-i", video_path,
                "-map", "0:s:0",  # First subtitle stream
                "-f", "srt",
                "-",  # Output to stdout
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )

        if result.returncode == 0 and result.stdout.strip():
            # Parse SRT and extract just the text (no timestamps)
            lines = result.stdout.strip().split("\n")
            text_lines = []
            for line in lines:
                line = line.strip()
                # Skip sequence numbers, timestamps, and empty lines
                if not line:
                    continue
                if line.isdigit():
                    continue
                if "-->" in line:
                    continue
                # Remove HTML tags from subtitles
                clean = re.sub(r"<[^>]+>", "", line)
                if clean.strip():
                    text_lines.append(clean.strip())

            return "\n".join(text_lines)
    except (subprocess.TimeoutExpired, FileNotFoundError):
        pass

    return ""


def _parse_subtitle_file(file_path: str) -> ParseResult:
    """Parse standalone subtitle files (SRT, ASS, SSA, VTT)."""
    try:
        if SUBS_AVAILABLE:
            subs = pysubs2.load(file_path)
            text_lines = []
            for event in subs:
                if event.text.strip():
                    # Clean subtitle formatting tags
                    clean = re.sub(r"\{[^}]*\}", "", event.text)
                    clean = re.sub(r"<[^>]+>", "", clean)
                    clean = clean.replace("\\N", "\n").replace("\\n", "\n").strip()
                    if clean:
                        text_lines.append(clean)

            subtitle_text = "\n".join(text_lines)
        else:
            # Fallback: read as plain text
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                subtitle_text = f.read()

        if not subtitle_text.strip():
            return ParseResult(
                pages=[PageContent(page_num=1, text="(Empty subtitle file)", content_type="subtitle")],
                metadata={"format": "subtitle", "page_count": 1, "source": file_path},
            )

        return ParseResult(
            pages=[PageContent(
                page_num=1,
                text=f"[SUBTITLES: {os.path.basename(file_path)}]\n{subtitle_text}",
                content_type="subtitle",
            )],
            metadata={
                "format": "subtitle",
                "page_count": 1,
                "source": file_path,
            },
        )
    except Exception as e:
        return ParseResult(success=False, error=f"Subtitle parse error: {e}")


# ---------------------------------------------------------------------------
# Code Parser
# ---------------------------------------------------------------------------
def _parse_code(file_path: str) -> ParseResult:
    """Parse code files and explicitly preserve formatting."""
    try:
        encoding = "utf-8"
        if CHARDET_AVAILABLE:
            with open(file_path, "rb") as f:
                raw = f.read(10000)
                detected = chardet.detect(raw)
                encoding = detected.get("encoding", "utf-8") or "utf-8"

        with open(file_path, "r", encoding=encoding, errors="replace") as f:
            content = f.read()

        if not content.strip():
            return ParseResult(
                pages=[PageContent(page_num=1, text="(Empty code file)")],
                metadata={"format": "code", "page_count": 1, "source": file_path},
            )

        ext = os.path.splitext(file_path)[1][1:]
        formatted_code = f"[CODE FILE: {os.path.basename(file_path)}]\n```{ext}\n{content}\n```"

        return ParseResult(
            pages=[PageContent(
                page_num=1,
                text=formatted_code,
            )],
            metadata={"format": "code", "page_count": 1, "source": file_path},
        )
    except Exception as e:
        return ParseResult(success=False, error=f"Code parse error: {e}")


# ---------------------------------------------------------------------------
# Email Parser
# ---------------------------------------------------------------------------
def _parse_email(file_path: str) -> ParseResult:
    """Parse .eml / .msg files into structured text."""
    try:
        with open(file_path, "rb") as f:
            msg = email.message_from_binary_file(f, policy=policy.default)

        subject = msg.get("subject", "(No Subject)")
        sender = msg.get("from", "(Unknown Sender)")
        date = msg.get("date", "(Unknown Date)")
        to = msg.get("to", "(Unknown Recipient)")

        body = ""
        # Extract body
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    content = part.get_content()
                    if content:
                        body += content + "\n"
        else:
            content = msg.get_content()
            if content:
                body = content

        if not body.strip() and msg.is_multipart():
            # Fallback to html if no plain text
            for part in msg.walk():
                if part.get_content_type() == "text/html":
                    html = part.get_content()
                    if html:
                        if BS4_AVAILABLE:
                            soup = BeautifulSoup(html, "html.parser")
                            body += soup.get_text() + "\n"
                        else:
                            body += re.sub(r'<[^>]+>', '', html) + "\n"

        email_text = (
            f"[EMAIL MESSAGE]\n"
            f"Subject: {subject}\n"
            f"From: {sender}\n"
            f"To: {to}\n"
            f"Date: {date}\n"
            f"---\n{body.strip()}"
        )

        return ParseResult(
            pages=[PageContent(
                page_num=1,
                text=email_text,
            )],
            metadata={"format": "email", "page_count": 1, "source": file_path},
        )
    except Exception as e:
        return ParseResult(success=False, error=f"Email parse error: {e}")


# ---------------------------------------------------------------------------
# URL Link Parser
# ---------------------------------------------------------------------------
def _parse_url(file_path: str) -> ParseResult:
    """Parse .url or .webloc and scrape offline content."""
    try:
        url = None
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
            # Find HTTP links
            match = re.search(r'(https?://[^\s]+)', content)
            if match:
                url = match.group(1)
                # Cleanup common trailing characters from markdown/XML context
                url = re.sub(r'["\]<].*$', '', url)

        if not url:
            return ParseResult(success=False, error="No valid URL found in file")

        if os.getenv("RAG_HF_OFFLINE", "false").lower() in {"1", "true", "yes", "on"}:
            return ParseResult(
                pages=[PageContent(page_num=1, text=f"[WEB URL]\n{url}\n\n(Offline mode: content not scraped)")],
                metadata={"format": "url", "page_count": 1, "source": url},
            )

        # Scrape content offline
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=10) as response:
            html = response.read().decode('utf-8', errors='replace')

        title = "Web Page"
        body_text = html
        if BS4_AVAILABLE:
            soup = BeautifulSoup(html, "html.parser")
            if soup.title:
                title = soup.title.get_text(strip=True) or "Web Page"
            # Remove scripts and styles
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()
            body_text = soup.get_text(separator="\n", strip=True)

        web_content = f"[WEB PAGE]\nTitle: {title}\nURL: {url}\n\n{body_text}"
        
        # Paginate very long web pages
        pages = []
        lines = web_content.split("\n")
        page_size = 100
        for i in range(0, max(len(lines), 1), page_size):
            chunk = lines[i:i + page_size]
            pages.append(PageContent(page_num=(i//page_size)+1, text="\n".join(chunk)))

        return ParseResult(
            pages=pages,
            metadata={"format": "url", "page_count": len(pages), "source": url},
        )
    except Exception as e:
        return ParseResult(success=False, error=f"URL scrape error: {e}")


# ---------------------------------------------------------------------------
# Archive Parser (zip, tar, gz, tgz, rar)
# ---------------------------------------------------------------------------
def _parse_archive(file_path: str) -> ParseResult:
    """Extract supported files from archives and parse each one."""
    import zipfile, tarfile

    pages = []
    supported_files = []
    extract_dir = None
    total_extracted = 0

    try:
        extract_dir = file_path + "_extracted"
        os.makedirs(extract_dir, exist_ok=True)
        extract_dir = os.path.realpath(extract_dir)

        def _safe_extract_zip(member_name: str) -> bool:
            """Extract a ZIP member safely; returns True if extracted, False if skipped."""
            dest = os.path.realpath(os.path.join(extract_dir, member_name))
            if not dest.startswith(extract_dir + os.sep):
                print(f"[Parser] Skipping ZIP member with path traversal: {member_name}")
                return False
            return True

        def _safe_extract_tar(member) -> bool:
            """Extract a TAR member safely; returns True if extracted, False if skipped."""
            if member.name.startswith(os.sep):
                print(f"[Parser] Skipping TAR member with absolute path: {member.name}")
                return False
            dest = os.path.realpath(os.path.join(extract_dir, member.name))
            if not dest.startswith(extract_dir + os.sep):
                print(f"[Parser] Skipping TAR member with path traversal: {member.name}")
                return False
            return True

        if file_path.endswith('.zip'):
            with zipfile.ZipFile(file_path, 'r') as zf:
                members = [m for m in zf.namelist()
                          if os.path.splitext(m)[1].lower() in SUPPORTED_EXTENSIONS
                          and not os.path.basename(m).startswith('.')]
                for member in members:
                    info = zf.getinfo(member)
                    total_extracted += info.file_size
                    if total_extracted > MAX_ARCHIVE_EXTRACT_SIZE:
                        print(f"[Parser] Archive extract aborted: exceeded {MAX_ARCHIVE_EXTRACT_SIZE} bytes")
                        break
                    if not _safe_extract_zip(member):
                        continue
                    zf.extract(member, extract_dir)
        elif file_path.endswith(('.tar', '.tar.gz', '.tgz')):
            with tarfile.open(file_path, 'r:*') as tf:
                members = [m for m in tf.getmembers()
                          if m.isfile()
                          and os.path.splitext(m.name)[1].lower() in SUPPORTED_EXTENSIONS
                          and not os.path.basename(m.name).startswith('.')]
                for member in members:
                    total_extracted += member.size
                    if total_extracted > MAX_ARCHIVE_EXTRACT_SIZE:
                        print(f"[Parser] Archive extract aborted: exceeded {MAX_ARCHIVE_EXTRACT_SIZE} bytes")
                        break
                    if not _safe_extract_tar(member):
                        total_extracted -= member.size
                        continue
                    tf.extract(member, extract_dir)

        for root, dirs, files in os.walk(extract_dir):
            for fname in files:
                fpath = os.path.join(root, fname)
                ext = os.path.splitext(fname)[1].lower()
                if ext in SUPPORTED_EXTENSIONS and ext not in ARCHIVE_EXTENSIONS:
                    supported_files.append(fpath)

        if not supported_files:
            return ParseResult(
                pages=[PageContent(page_num=1, text="(Archive contained no supported files)")],
                metadata={"format": "archive", "page_count": 1, "source": file_path},
            )

        for fpath in supported_files:
            result = parse_file(fpath)
            if result.success:
                pages.extend(result.pages)
            else:
                print(f"[Parser] Archive sub-file failed: {fpath}: {result.error}")

        return ParseResult(
            pages=pages,
            metadata={
                "format": "archive",
                "page_count": len(pages),
                "source": file_path,
                "contained_files": len(supported_files),
            },
        )
    finally:
        if extract_dir and os.path.isdir(extract_dir):
            shutil.rmtree(extract_dir, ignore_errors=True)


# ---------------------------------------------------------------------------
# Fallback parser wrappers — try Docling/MinerU first, fall back to native
# ---------------------------------------------------------------------------
def _fallback_parse_pdf(file_path: str) -> ParseResult:
    """Fallback PDF parser using pdfplumber."""
    if not PDFPLUMBER_AVAILABLE:
        return ParseResult(success=False, error="pdfplumber not installed for PDF fallback")
    try:
        pages = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                tables = [str(t) for t in page.extract_tables()] if page.extract_tables() else []
                pages.append(PageContent(page_num=i + 1, text=text, tables=tables))
        return ParseResult(
            pages=pages,
            metadata={"format": "pdf", "page_count": len(pages), "source": file_path, "parser": "fallback_pdfplumber"},
        )
    except Exception as e:
        return ParseResult(success=False, error=f"PDF fallback failed: {e}")

def _fallback_parse_docx(file_path: str) -> ParseResult:
    """Fallback DOCX parser using python-docx."""
    if not DOCX_AVAILABLE:
        return ParseResult(success=False, error="python-docx not installed for DOCX fallback")
    try:
        doc = Document(file_path)
        text = "\n".join(p.text for p in doc.paragraphs)
        pages = [PageContent(page_num=1, text=text)]
        return ParseResult(
            pages=pages,
            metadata={"format": "docx", "page_count": 1, "source": file_path, "parser": "fallback_docx"},
        )
    except Exception as e:
        return ParseResult(success=False, error=f"DOCX fallback failed: {e}")

def _fallback_parse_xlsx(file_path: str) -> ParseResult:
    """Fallback XLSX parser using openpyxl."""
    if not OPENPYXL_AVAILABLE:
        return ParseResult(success=False, error="openpyxl not installed for XLSX fallback")
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        try:
            pages = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(c) if c is not None else "" for c in row)
                    rows.append(row_text)
                text = f"[Sheet: {sheet_name}]\n" + "\n".join(rows)
                pages.append(PageContent(page_num=len(pages) + 1, text=text))
        finally:
            wb.close()
        return ParseResult(
            pages=pages,
            metadata={"format": "xlsx", "page_count": len(pages), "source": file_path, "parser": "fallback_openpyxl"},
        )
    except Exception as e:
        return ParseResult(success=False, error=f"XLSX fallback failed: {e}")

def _fallback_parse_pptx(file_path: str) -> ParseResult:
    """Fallback PPTX parser using python-pptx."""
    if not PPTX_AVAILABLE:
        return ParseResult(success=False, error="python-pptx not installed for PPTX fallback")
    try:
        prs = Presentation(file_path)
        pages = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
            pages.append(PageContent(page_num=i + 1, text="\n".join(texts)))
        return ParseResult(
            pages=pages,
            metadata={"format": "pptx", "page_count": len(pages), "source": file_path, "parser": "fallback_pptx"},
        )
    except Exception as e:
        return ParseResult(success=False, error=f"PPTX fallback failed: {e}")

def _fallback_parse_html(file_path: str) -> ParseResult:
    """Fallback HTML parser using BeautifulSoup."""
    if not BS4_AVAILABLE:
        return ParseResult(success=False, error="BeautifulSoup not installed for HTML fallback")
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        pages = [PageContent(page_num=1, text=text)]
        return ParseResult(
            pages=pages,
            metadata={"format": "html", "page_count": 1, "source": file_path, "parser": "fallback_bs4"},
        )
    except Exception as e:
        return ParseResult(success=False, error=f"HTML fallback failed: {e}")

def _parse_pdf_with_fallback(file_path: str) -> ParseResult:
    """Parse PDF: Docling → pdfplumber."""
    result = _parse_docling(file_path)
    if result.success:
        return result
    print(f"[Parser] Docling failed, falling back to pdfplumber: {result.error}")
    return _fallback_parse_pdf(file_path)

def _parse_office_with_fallback(file_path: str) -> ParseResult:
    """Parse Office: Docling first, fall back to format-specific native parser."""
    result = _parse_docling(file_path)
    if result.success:
        return result
    print(f"[Parser] Docling failed, falling back to native parser: {result.error}")
    ext = os.path.splitext(file_path)[1].lower()
    native = {
        ".docx": _fallback_parse_docx, ".doc": _fallback_parse_docx,
        ".xlsx": _fallback_parse_xlsx, ".xls": _fallback_parse_xlsx,
        ".pptx": _fallback_parse_pptx, ".ppt": _fallback_parse_pptx,
    }.get(ext)
    if native:
        fb = native(file_path)
        if fb.success:
            return fb
    return result

def _parse_html_with_fallback(file_path: str) -> ParseResult:
    """Parse HTML: Docling first, fall back to BeautifulSoup."""
    result = _parse_docling(file_path)
    if result.success:
        return result
    print(f"[Parser] Docling failed for HTML, falling back to BeautifulSoup: {result.error}")
    return _fallback_parse_html(file_path)


# ---------------------------------------------------------------------------
# Main Entry Point — Universal Parser
# ---------------------------------------------------------------------------
PARSER_REGISTRY = {
    ".pdf": _parse_pdf_with_fallback,
    ".docx": _parse_office_with_fallback,
    ".doc": _parse_office_with_fallback,
    ".xlsx": _parse_office_with_fallback,
    ".xls": _parse_office_with_fallback,
    ".pptx": _parse_office_with_fallback,
    ".ppt": _parse_office_with_fallback,
    ".html": _parse_html_with_fallback,
    ".htm": _parse_html_with_fallback,
    ".csv": _parse_csv,
    ".txt": _parse_text,
    ".text": _parse_text,
    ".md": _parse_text,
    ".log": _parse_text,
    ".json": _parse_text,
    ".xml": _parse_text,
    # Images
    ".png": _parse_image,
    ".jpg": _parse_image,
    ".jpeg": _parse_image,
    ".bmp": _parse_image,
    ".tiff": _parse_image,
    ".tif": _parse_image,
    ".gif": _parse_image,
    ".webp": _parse_image,
    # Video
    ".mp4": _parse_video,
    ".avi": _parse_video,
    ".mkv": _parse_video,
    ".mov": _parse_video,
    ".wmv": _parse_video,
    ".flv": _parse_video,
    # Subtitle files
    ".srt": _parse_subtitle_file,
    ".ass": _parse_subtitle_file,
    ".ssa": _parse_subtitle_file,
    ".vtt": _parse_subtitle_file,
}

# Add code extensions dynamically
for ext in CODE_EXTENSIONS:
    PARSER_REGISTRY[ext] = _parse_code

# Add email extensions dynamically
for ext in EMAIL_EXTENSIONS:
    PARSER_REGISTRY[ext] = _parse_email

# Add URL extensions dynamically
for ext in URL_EXTENSIONS:
    PARSER_REGISTRY[ext] = _parse_url

# Add archive extensions dynamically
for ext in ARCHIVE_EXTENSIONS:
    PARSER_REGISTRY[ext] = _parse_archive


# ---------------------------------------------------------------------------
# Fallback parser wrappers — try Docling/MinerU first, fall back to native

def parse_file(file_path: str) -> ParseResult:
    """
    Universal file parser — auto-detects format and extracts all content.
    Uses fallback chains: Docling/MinerU → native format-specific parsers.
    """
    if not os.path.exists(file_path):
        return ParseResult(success=False, error=f"File not found: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()
    parser = PARSER_REGISTRY.get(ext)

    if not parser:
        return ParseResult(
            success=False,
            error=f"Unsupported file format: {ext}. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}",
        )

    print(f"[Parser] Parsing {os.path.basename(file_path)} (format: {ext})")
    result = parser(file_path)

    if result.success:
        total_text = sum(len(p.text) for p in result.pages)
        total_tables = sum(len(p.tables) for p in result.pages)
        total_images = sum(len(p.image_texts) for p in result.pages)
        print(
            f"[Parser] ✅ Parsed {len(result.pages)} pages, "
            f"{total_text} chars, {total_tables} tables, {total_images} images"
        )
    else:
        print(f"[Parser] ❌ Failed: {result.error}")

    return result
